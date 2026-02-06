"""PPTX to Markdown converter for PowerPoint presentations."""

from pathlib import Path
from typing import Any

from fileconverter.converters.base import BaseConverter, ConversionResult


class PPTXConverter(BaseConverter):
    """Convert PPTX (PowerPoint) files to Markdown."""

    @classmethod
    def supported_extensions(cls) -> list[str]:
        return [".pptx"]

    def convert(self, file_path: Path) -> ConversionResult:
        """Convert a PPTX file to Markdown.

        Args:
            file_path: Path to the PPTX file.

        Returns:
            ConversionResult with the conversion outcome.
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as e:
            return self._create_error_result(
                file_path,
                f"Required package not installed: {e}. Run: pip install python-pptx",
            )

        try:
            # Open the presentation
            prs = Presentation(str(file_path))

            markdown_parts = []
            images_list: list[Path] = []

            # Add title
            markdown_parts.append(f"# {file_path.stem}\n")

            # Setup image directory if needed
            images_dir = None
            if self.extract_images:
                images_dir = file_path.parent / f"{file_path.stem}_images"
                images_dir.mkdir(exist_ok=True)

            image_counter = 0

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, start=1):
                markdown_parts.append(f"\n---\n\n## Slide {slide_num}")

                # Extract slide title if it exists
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        markdown_parts.append(f": {title_text}")

                markdown_parts.append("\n\n")

                # Extract content from shapes
                for shape in slide.shapes:
                    # Skip the title shape (already processed)
                    if shape == slide.shapes.title:
                        continue

                    # Handle text boxes and shapes with text
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # Check if it's a text frame with paragraphs
                        if hasattr(shape, "text_frame"):
                            formatted_text = self._extract_text_frame(shape.text_frame)
                            if formatted_text:
                                markdown_parts.append(formatted_text + "\n\n")
                        else:
                            markdown_parts.append(text + "\n\n")

                    # Handle tables
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_md = self._extract_table(shape.table)
                        if table_md:
                            markdown_parts.append(table_md + "\n\n")

                    # Handle pictures
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if self.extract_images and images_dir:
                            try:
                                image_counter += 1
                                image = shape.image

                                # Get image extension from content type
                                ext = image.ext
                                image_filename = f"slide_{slide_num}_image_{image_counter}.{ext}"
                                image_path = images_dir / image_filename

                                # Save the image
                                with open(image_path, 'wb') as f:
                                    f.write(image.blob)

                                images_list.append(image_path)

                                # Add image reference to markdown
                                markdown_parts.append(f"![Image {image_counter}]({image_path})\n\n")
                            except Exception as e:
                                markdown_parts.append(f"*[Image extraction failed: {str(e)}]*\n\n")
                        else:
                            markdown_parts.append(f"*[Image: slide_{slide_num}_image_{image_counter + 1}]*\n\n")
                            image_counter += 1

                    # Handle groups (recursively process grouped shapes)
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        group_text = self._extract_group_text(shape)
                        if group_text:
                            markdown_parts.append(group_text + "\n\n")

                # Extract speaker notes if they exist
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        markdown_parts.append(f"**Speaker Notes:**\n\n{notes_text}\n\n")

            # Combine all parts
            markdown = "".join(markdown_parts).strip()

            # Clean up the markdown
            markdown = self._clean_markdown(markdown)

            return self._create_success_result(file_path, markdown, images_list)

        except Exception as e:
            return self._create_error_result(file_path, f"Failed to convert PPTX: {str(e)}")

    def _extract_text_frame(self, text_frame: Any) -> str:
        """Extract formatted text from a text frame.

        Args:
            text_frame: PowerPoint text frame object.

        Returns:
            Formatted markdown text.
        """
        lines = []

        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            # Check indentation level for bullets
            level = paragraph.level

            if level > 0:
                # Indented bullet point
                indent = "  " * level
                lines.append(f"{indent}- {text}")
            else:
                # Top level
                lines.append(f"- {text}")

        return "\n".join(lines) if lines else ""

    def _extract_table(self, table: Any) -> str:
        """Extract table and convert to markdown.

        Args:
            table: PowerPoint table object.

        Returns:
            Markdown formatted table.
        """
        try:
            rows = []

            for row in table.rows:
                cells = []
                for cell in row.cells:
                    # Get cell text and clean it
                    cell_text = cell.text.strip().replace('\n', ' ').replace('|', '\\|')
                    cells.append(cell_text)
                rows.append(cells)

            if not rows:
                return ""

            # Calculate column widths
            col_count = len(rows[0])
            col_widths = [max(len(row[i]) if i < len(row) else 0 for row in rows) for i in range(col_count)]
            col_widths = [max(3, w) for w in col_widths]

            markdown_lines = []

            # First row as header
            if rows:
                header = rows[0]
                header_cells = [cell.ljust(col_widths[i]) for i, cell in enumerate(header)]
                markdown_lines.append("| " + " | ".join(header_cells) + " |")

                # Separator
                separators = ["-" * width for width in col_widths]
                markdown_lines.append("| " + " | ".join(separators) + " |")

                # Data rows
                for row in rows[1:]:
                    cells = [row[i].ljust(col_widths[i]) if i < len(row) else " " * col_widths[i]
                             for i in range(col_count)]
                    markdown_lines.append("| " + " | ".join(cells) + " |")

            return "\n".join(markdown_lines)

        except Exception:
            return ""

    def _extract_group_text(self, group_shape: Any) -> str:
        """Extract text from grouped shapes.

        Args:
            group_shape: PowerPoint group shape object.

        Returns:
            Combined text from all shapes in the group.
        """
        texts = []

        for shape in group_shape.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        return "\n".join(texts) if texts else ""

    def _clean_markdown(self, markdown: str) -> str:
        """Clean up the markdown output.

        Args:
            markdown: Raw markdown text.

        Returns:
            Cleaned markdown text.
        """
        # Remove excessive blank lines (more than 2 consecutive)
        lines = markdown.split("\n")
        cleaned_lines: list[str] = []
        blank_count = 0

        for line in lines:
            is_blank = not line.strip()
            if is_blank:
                blank_count += 1
                if blank_count <= 2:
                    cleaned_lines.append(line)
            else:
                blank_count = 0
                cleaned_lines.append(line)

        return "\n".join(cleaned_lines).strip()
